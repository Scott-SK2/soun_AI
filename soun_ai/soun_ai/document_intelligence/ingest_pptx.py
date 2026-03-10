from __future__ import annotations
from typing import List, Optional
import os
import re

from pptx import Presentation
from document_intelligence.document_store import DocChunk


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _chunk_text(text: str, max_chars: int = 800, overlap: int = 120) -> List[str]:
    text = _clean(text)
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        j = min(len(text), i + max_chars)
        chunk = text[i:j].strip()
        if chunk:
            chunks.append(chunk)
        if j == len(text):
            break
        i = max(0, j - overlap)
    return chunks


def ingest_pptx(path: str, source_name: Optional[str] = None) -> List[DocChunk]:
    source = source_name or os.path.basename(path)
    prs = Presentation(path)

    out: List[DocChunk] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # --- Slide text ---
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text)

        slide_joined = _clean("\n".join(slide_texts))
        if slide_joined:
            for ch in _chunk_text(slide_joined):
                out.append(DocChunk(text=ch, source=source, page=slide_idx, kind="speaker_notes"))

        # --- Speaker notes (notes pane) ---
        try:
            if slide.has_notes_slide and slide.notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text:
                    notes_text = _clean(notes_frame.text)
                    if notes_text:
                        for ch in _chunk_text(notes_text):
                            # page = slide index; keep same mapping
                            out.append(DocChunk(text=ch, source=source, page=slide_idx, kind="speaker_notes"))
        except Exception:
            # Notes can be missing or malformed; ignore safely
            pass

    return out